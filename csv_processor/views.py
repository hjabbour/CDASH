
# Create your views here.
# 2023 Haytham Jabbour hjabbour
from django.shortcuts import render, redirect
from pymongo import MongoClient
import pandas as pd
from .forms import UploadCSVForm
#from .forms import UploadCSVForm
from .models import UploadedFile
from django.conf import settings
from django.http import HttpResponse
from pptx import Presentation
import json
from bson import json_util


#from pandas_datareader import data
from datetime import date
client = MongoClient('mongodb://root:rootpassword@192.168.2.152:27017')
db = client['CDASH']
collection = db['csvdash']

import os


def delete_files(request):
    if request.method == 'POST':
        selected_files = request.POST.getlist('selected_files')
        if selected_files:
            # Delete files from the file system
            folder_path = os.path.join(settings.MEDIA_ROOT, 'csv_files')
            for file_name in selected_files:
                file_path = os.path.join(folder_path, file_name)
                if os.path.exists(file_path):
                    os.remove(file_path)

            # Delete records from the database
            UploadedFile.objects.filter(csv_file__in=selected_files).delete()
            
        return redirect('csv_processor:file_list')
    return render(request, 'csv_processor/file_list.html')



def upload_csv(request):
    if request.method == 'POST':
        form = UploadCSVForm(request.POST, request.FILES)
        if form.is_valid():
            csv_file = request.FILES['csv_file']

            # Save the uploaded file
            uploaded_file = UploadedFile(csv_file=csv_file)
            uploaded_file.save()

            # Redirect to the file list view on successful upload
            return redirect('csv_processor:file_list')
    else:
        form = UploadCSVForm()

    context = {
        'form': form
    }
    return render(request, 'csv_processor/upload_csv.html', context)


def file_list(request):
    # Get the list of files in the csv_files folder
    folder_path = os.path.join(settings.MEDIA_ROOT, 'csv_files')
    files = os.listdir(folder_path)
    
    context = {
        'files': files
    }
    return render(request, 'csv_processor/file_list.html', context)

def file_detail(request, filename):
    # Open and read the selected file
    file_path = os.path.join(settings.MEDIA_ROOT, 'csv_files', filename)
    with open(file_path, 'r') as file:
        content = file.read()
    
    context = {
        'filename': filename,
        'content': content
    }
    return render(request, 'csv_processor/file_detail.html', context)



def file_columns(request, filename):
    file_path = os.path.join(settings.MEDIA_ROOT, 'csv_files', filename)
    # Code to read the CSV file and extract the column names
    # You can use the csv module or pandas to read the file and extract the columns
    
    # Example using pandas:
    import pandas as pd
    df = pd.read_csv(file_path)
    columns = df.columns.tolist()
    
    context = {
        'filename': filename,
        'selected_file': filename,
        'columns': columns,
    
    }
    return render(request, 'csv_processor/file_columns.html', context)


def process_columns(request,filename):
    if request.method == 'POST':
        
        # Retrieve selected columns from the form data
        selected_columns = request.POST.getlist('columns')
        selected_file = request.POST.get('filename')
        file_path = os.path.join(settings.MEDIA_ROOT, 'csv_files', filename)
        df = pd.read_csv(file_path)
        new_df=df[selected_columns]
        
        context = {
        'filename': filename,
        'selected_file': filename,
        'columns': selected_columns,
        'df':new_df,
    
    }
        
    return render(request, 'csv_processor/new_columns.html', context)

        # Process the selected columns
        # ...

        # Return an appropriate response, such as a redirect or a success message


    # Process the selected columns as needed (insert into MongoDB or return as pandas dataframe)
    #
def process_ppt(request):
    if request.method == 'POST':
        ppt_file = request.FILES['ppt_file']

        # Read the uploaded PPT file
        presentation = Presentation(ppt_file)

        # Extract slide contents
        slides = []
        for slide in presentation.slides:
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_content.append(run.text)
            slides.append(slide_content)

        # Pass the extracted data to the template
        context = {
            'slides': slides
        }
        return render(request, 'ppt_processor/result.html', context)

    return render(request, 'ppt_processor/upload_ppt.html')

## fix this its not inserting 
def insert_into_mongodb(request):
    if request.method == 'POST':
        filename = request.POST['filename']
        selected_columns = request.POST.getlist('columns')
        file_path = os.path.join(settings.MEDIA_ROOT, 'csv_files', filename)
        df = pd.read_csv(file_path)
        new_df = df[selected_columns]
        
        
        client = MongoClient('mongodb://root:password@192.168.2.152:27017')
        db = client['CDASH']
        collection = db['csvdash']
        # Convert DataFrame to JSON
        #df_json = df.to_json(orient='records')
        
        # Connect to MongoDB
        #client = MongoClient('mongodb://localhost:27017/')
        #db = client['your_database_name']
        #collection = db['your_collection_name']
        
        # Insert JSON data into MongoDB
        #json_data = json.loads(df_json)
        
        data_dict = df.to_dict("records")
        result=collection.insert_many(data_dict)
      
        if result:
             return render(request,'csv_processor/success.html')
        else:
             return render(request,'csv_processor/error.html')
        #client.close()
        
        # Redirect to a success page or another URL
        return render(request,'csv_processor/success.html')
    
    # Handle GET request or invalid form submission
    return render(request,'csv_processor/error.html')